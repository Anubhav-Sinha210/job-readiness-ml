from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()

    # ---------------------------
    # Slide 1: Title Slide
    # ---------------------------
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Student Job Readiness Prediction System"
    subtitle.text = "End-to-End Machine Learning Project\nPrepared for the CSE274 Project"

    # ---------------------------
    # Slide 2: Project Overview
    # ---------------------------
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Project Overview & Objectives"
    tf = body_shape.text_frame
    tf.text = "Problem Statement: Predict whether a student is job-ready based on academic and extracurricular metrics."
    p = tf.add_paragraph()
    p.text = "Objective: Provide an automated evaluation system with actionable insights and personalized recommendations for students."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Technology Stack:"
    p = tf.add_paragraph()
    p.text = "Python (Data Processing & Modeling)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Scikit-learn, XGBoost, imbalanced-learn (Machine Learning)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Streamlit (Web Interface)"
    p.level = 1

    # ---------------------------
    # Slide 3: Data Attributes
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Key Data Attributes"
    tf = body_shape.text_frame
    tf.text = "Academic Metrics:"
    p = tf.add_paragraph()
    p.text = "Degree Stream, CGPA"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Skill Assessments (Out of 100):"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Technical Skills Score, Aptitude Score"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Communication Skills, Coding Score"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Practical Experience:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Number of Internships"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Number of Projects"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Number of Certifications"
    p.level = 1

    # ---------------------------
    # Slide 4: Data Preprocessing
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Data Preprocessing Pipeline"
    tf = body_shape.text_frame
    tf.text = "Handling Data Issues: Cleaning missing values and extreme anomalies."
    p = tf.add_paragraph()
    p.text = "Outlier Treatment: Employed the Interquartile Range (IQR) method."
    p = tf.add_paragraph()
    p.text = "Standardization: Applied StandardScaler to normalize scoring distributions."
    p = tf.add_paragraph()
    p.text = "Categorical Encoding: One-Hot-Encoding for Stream degrees."
    p = tf.add_paragraph()
    p.text = "Handling Class Imbalance: Applied SMOTE (Synthetic Minority Over-sampling Technique) to ensure the model doesn't bias towards the majority class."

    # ---------------------------
    # Slide 5: Machine Learning Models
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Modeling & Evaluation"
    tf = body_shape.text_frame
    tf.text = "Supervised Models Evaluated:"
    p = tf.add_paragraph()
    p.text = "Logistic Regression, Naive Bayes, KNN, SVM, Decision Tree"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Advanced Ensemble Models:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Random Forest, XGBoost"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Voting Classifier aggregating top performing models"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Training Strategies:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Hyperparameter Tuning using GridSearchCV"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Feature Selection relying on Feature Importance and Variance Threshold"
    p.level = 1

    # ---------------------------
    # Slide 6: Insights & App Backend
    # ---------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Unsupervised Insights & Production"
    tf = body_shape.text_frame
    tf.text = "Clustering Analysis:"
    p = tf.add_paragraph()
    p.text = "K-Means Clustering utilized to identify distinct student talent cohorts."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Optimal clusters determined via Elbow and Silhouette Plots."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Interactive Application (Streamlit):"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Uses saved inference pipelines (preprocessor -> selector -> best_model)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Outputs visual readiness probability percentages and dynamic recommendations based on input profile."
    p.level = 1

    # Save presentation
    output_path = "Student_Job_Readiness_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == '__main__':
    create_presentation()
